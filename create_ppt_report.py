from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.oxml.xmlchemy import OxmlElement
from google.oauth2 import service_account
from googleapiclient.discovery import build
from googleapiclient.http import MediaIoBaseDownload
from oauth2client.service_account import ServiceAccountCredentials
from datetime import datetime

import os


def SubElement(parent, tagname, **kwargs):
    element = OxmlElement(tagname)
    element.attrib.update(kwargs)
    parent.append(element)
    return element

def _set_cell_border(cell, border_color="000000", border_width='12700'):
    tc = cell._tc
    tcPr = tc.get_or_add_tcPr()
    for lines in ['a:lnL','a:lnR','a:lnT','a:lnB']:

        # Every time before a node is inserted, the nodes with the same tag should be removed.
        tag = lines.split(":")[-1]
        for e in tcPr.getchildren():
            if tag in str(e.tag):
                tcPr.remove(e)
        # end

        ln = SubElement(tcPr, lines, w=border_width, cap='flat', cmpd='sng', algn='ctr')
        solidFill = SubElement(ln, 'a:solidFill')
        srgbClr = SubElement(solidFill, 'a:srgbClr', val=border_color)
        prstDash = SubElement(ln, 'a:prstDash', val='solid')
        round_ = SubElement(ln, 'a:round')
        headEnd = SubElement(ln, 'a:headEnd', type='none', w='med', len='med')
        tailEnd = SubElement(ln, 'a:tailEnd', type='none', w='med', len='med')

# Function to download the image from a URL
def download_image(service_account_file, file_path, save_path):
    # Authenticate using service account credentials
    credentials = service_account.Credentials.from_service_account_file(service_account_file)
    drive_service = build('drive', 'v3', credentials=credentials, cache_discovery=False)
    file_id = get_file_id_from_url(file_path)
    # Download the image
    request = drive_service.files().get_media(fileId=file_id)
    fh = open(save_path, "wb")
    downloader = MediaIoBaseDownload(fh, request)
    done = False
    while done is False:
        status, done = downloader.next_chunk()
        print("Download %d%%." % int(status.progress() * 100))
    print("Image downloaded successfully.")

def get_file_id_from_url(url):
    # Extract the file ID from the URL
    return url.split('/')[-2]

############################ Google Sheets ##########################################
import gspread

# Define the scope and credentials for Google Sheets API
scope = ['https://spreadsheets.google.com/feeds',
   'https://www.googleapis.com/auth/drive']

# Authorize the client with credentials
creds = ServiceAccountCredentials.from_json_keyfile_name('/home/sailbslsafety/mysite/credentials.json', scope)
client = gspread.authorize(creds)

######################################################################################

def create_slide(ppt, insp_category, observation):
    print ("in create_slide")
    blank_slide_layout = ppt.slide_layouts[6]
    slide = ppt.slides.add_slide(blank_slide_layout)
    shapes = slide.shapes
    left = Inches(0.1)
    top = Inches(5.5)
    width = Inches(9.7)
    height = Inches(0.8)
    if (insp_category=="General"):
        rows=4
    else:
        rows=5
    cols=4
    table = shapes.add_table(rows, cols, left, top, width, height).table

    for i in range (rows):
        row = table.rows[i].cells
        for j in range(cols):
            _set_cell_border(table.cell(i,j))

    #Readying the template
    if (insp_category=="General"):
        table.cell(0,0).text = "Department:"
        table.cell(0,2).text = "Location:"
        table.cell(1,0).text = "Observation:"
        table.cell(2,0).text = "Informed to:"
        table.cell(3,0).text = "Inspection Date:"
        table.cell(3,2).text = "Completion Target:"
        table.cell(1,1).merge(table.cell(1,3))
        table.cell(2,1).merge(table.cell(2,3))
    else:
        table.cell(0,0).text = "Activity"
        table.cell(0,1).text = observation['inspection_category']
        table.cell(1,0).text = "Department:"
        table.cell(1,2).text = "Location:"
        table.cell(2,0).text = "Observation:"
        table.cell(3,0).text = "Informed to:"
        table.cell(4,0).text = "Inspection Date:"
        table.cell(4,2).text = "Completion Target:"
        table.cell(0,1).merge(table.cell(0,3))
        table.cell(2,1).merge(table.cell(2,3))
        table.cell(3,1).merge(table.cell(3,3))

    #Feeding data to template
    if (insp_category=="General"):
        table.cell(0,1).text = observation['department']
        table.cell(0,3).text = observation['location']
        table.cell(1,1).text = observation['observation']
        table.cell(2,1).text = observation['discussed_with']
        table.cell(3,1).text = str(datetime.strptime(str(observation['date']), "%Y-%m-%d").strftime("%d-%m-%Y"))
        table.cell(3,3).text = observation['target_date']
        if (observation['compliance_status']=="Complied"):
            table.cell(3,3).text = "Complied"
    else:
        table.cell(1,1).text = observation['inspection_category']
        table.cell(1,1).text = observation['department']
        table.cell(1,3).text = observation['location']
        table.cell(2,1).text = observation['observation']
        table.cell(3,1).text = observation['discussed_with']
        table.cell(4,1).text = observation['date']
        table.cell(4,3).text = observation['target_date']
        if (observation['compliance_status']=="Complied"):
            table.cell(4,3).text = "Complied"
    if (observation['photo']!="" and observation['complied_photo']!=""):
        image_list = (observation['photo']+','+observation['complied_photo']).split(",")
    elif (observation['photo']!="" and observation['complied_photo']==""):
        image_list = observation['photo'].split(",")
    elif (observation['photo']=="" and observation['complied_photo']!=""):
        image_list = observation['complied_photo'].split(",")
    else:
        image_list=[]
    if (len(image_list)!=0):
        delta = 0   #used so that two or more images on a single slide do not completely overlap each other
        for image in image_list:
            download_image('/home/sailbslsafety/mysite/credentials.json', image, 'image.jpg')
            left = Inches(1+delta)  # Adjust the left position (in inches)
            top = Inches(1)  # Adjust the top position (in inches)
            width = Inches(2)  # Adjust the width (in inches)
            height = Inches(2)  # Adjust the height (in inches)
            slide.shapes.add_picture('image.jpg', left, top+delta, width, height)
            # Remove the image from local system
            os.remove('image.jpg')
            delta+=0.5


    # Add border, change font color to black and change background to white
    for i in range(rows):
        for j in range(cols):
            #set border
            _set_cell_border(table.cell(i,j))

            # Iterate through each paragraph in the cell
            for paragraph in table.cell(i,j).text_frame.paragraphs:
                # Iterate through each run in the paragraph and set font color to black
                for run in paragraph.runs:
                    run.font.color.rgb = RGBColor(0, 0, 0)  # Set font color to black

            # Set fill color for the cell
            table.cell(i,j).fill.solid()
            table.cell(i,j).fill.fore_color.rgb = RGBColor(255, 255, 255)  # White fill color

    #Give color (Red or Green) depending on compliance status of the observation.
    compliance_status = observation['compliance_status']
    if (insp_category=="General"):
        i=j=1
    else:
        i=2
        j=1
    if (compliance_status=="Good Point"):
        for paragraph in table.cell(i,j).text_frame.paragraphs:
            for run in paragraph.runs:
                run.font.color.rgb = RGBColor(0, 100, 0)
                run.font.bold=True
    else:
        for paragraph in table.cell(i,j).text_frame.paragraphs:
            for run in paragraph.runs:
                run.font.color.rgb = RGBColor(255, 0, 0)
                run.font.bold=True

def createPPT(inspection_df):
    print ("In createPPT")
    ppt = Presentation()
    for index, observation in inspection_df.iterrows():
        create_slide(ppt, observation['inspection_category'],observation)
    file_path = "/home/sailbslsafety/mysite/downloads" + "Presentation" +".pptx"
    ppt.save(file_path)
    return file_path